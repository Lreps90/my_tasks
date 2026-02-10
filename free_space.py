from pptx import Presentation
from pptx.util import Inches

# Initialize the presentation object
prs = Presentation()

# List of questions
questions = [
    """Question 1 (**)
    A uniform rod AB has length 5 m and weight 100 N. The rod rests in a horizontal position on two smooth supports at P and Q, where 1 mAP =, as shown in the figure above. The magnitude of the reaction force on the rod at P is 40 N.
    a) Determine the magnitude of the reaction force on the rod at Q.
    b) Calculate the distance AQ.""",

    """Question 2 (**)
    A non-uniform plank of wood AB has length 8 m and mass 100 kg. The plank is smoothly supported at its two ends A and B. A boy of mass 60 kg stands on the plank at the point C, where 3AC = m, as shown in the figure above. The plank with the boy standing on the plank, remains in equilibrium with AB horizontal. The plank is modelled as a non-uniform rod and the boy as a particle.
    a) Given that the reactions at the two supports are equal, determine the distance of the centre of mass of the plank from A.
    b) Explain in the context of this problem the model of
    i. … the plank is a rod
    ii. … the boy is a particle.""",

    """Question 3 (**)
    A plank of wood AB has length 4 m and mass 40 kg. The plank is smoothly supported at A and at C, where 3AC = m, as shown in the figure above. A man of mass 80 kg stands on the plank at a distance d m from A. The plank, with the man standing on it, remains in equilibrium with AB horizontal, and the reactions on the plank at A and at C equal. The plank is modelled as a uniform rod and the man as a particle.
    Determine the value of d.""",

    """Question 4 (**)
    A uniform iron girder AB has length 8 m and weight W N. A load of 250 N is attached to the girder at A and a load of 400 N is attached to the girder at B. The loaded girder is suspended by two light vertical cables attached to the girder at points C and D, where 1 mAC = and 3 mDB =. When the loaded girder rests undisturbed in a horizontal position, the tension in the cable at D is four times the tension at the cable at C. The girder is modelled as a uniform rod and the two loads as particles.
    a) Determine magnitude of the tension on the girder at C.
    b) Find the value of W.""",

    """Question 5 (**)
    A uniform rod AB has length 6 m and weight 40 kg. The rod rests in a horizontal position on two smooth supports at P and Q, where 1 mAP = and mAQ d= . The magnitude of the reaction force on the rod at Q is 3 times as large as that at P. Calculate the value of d.""",

    """Question 6 (**+)
    A box of mass 76 kg is attached by a string to one end B of a uniform rod AB of length 5 m and mass 24 kg. The rod is held horizontally in equilibrium by two smooth cylindrical pegs, one at A and one at C, where 2AC = m, as shown in the figure above. Calculate the magnitude of the forces exerted by each of the pegs onto the rod.""",

    """Question 7 (***)
    A beam AB has length 5.5 m and mass 20 kg. The beam is smoothly supported at the point P, where 2AP = m. A man of mass 70 kg stands on the beam at A and another man stands on the beam at a distance of 2.5 m from B. The beam is modelled as a non-uniform rod and the men are modelled as particles. The beam is in equilibrium in a horizontal position with the reaction on the beam at P having a magnitude of 1960 N. Calculate the distance of the centre of mass of the beam from A.""",

    """Question 8 (***)
    The figure above shows a uniform wooden beam AB, of length x m and weight 80 N. The beam is smoothly hinged at A and rests in a horizontal position on a smooth support at C, where 3 mAC =. When a rock of weight 70 N is placed on the beam at B the magnitude of the reaction force on the beam at C is 165 N. The beam is modelled as a uniform rod and the rock as a particle.
    a) Calculate the value of x.
    b) Explain briefly the model …
    i. … the beam is a uniform rod.
    ii. … the rock is a particle.
    The rock is next moved to a new position D on the beam, so that the beam with the rock at D remains in equilibrium in a horizontal position. The magnitude of reaction force at the support at C is now twenty times as large as the reaction force at the hinge at A.
    c) Calculate the distance AD.""",

    """Question 9 (***)
    A mechanical lever consists of a uniform steel rigid rod AB, of length 2 m and weight 100 N, placed over a smooth pivot at C. A box of weight 2400 N is suspended by a light inextensible string at B. When a vertical force is applied at A, as shown in the figure above, the lever remains in equilibrium, with AB horizontal.
    a) Given that 0.3CB = m, determine the magnitude of the force applied at A.
    The position of the pivot is changed so that lever remains in equilibrium when the vertical force applied at A has magnitude 200 N.
    b) Calculate the new distance of the pivot from B.""",

    """Question 10 (***)
    The figure above shows a uniform rod AB of length 1.8 m and mass 3 kg, held in a horizontal position by two small smooth pegs C and D. A particle of mass 12 kg, is placed at B. Given that 0.3AC = m and 0.4CD = m, determine the magnitude of each of the forces exerted on the rod by the pegs.""",

    """Question 11 (***)
    A non-uniform plank of wood AB has length 8.5 m and mass 20 kg. The centre of mass of the plank is 3.75 m from B. The plank is smoothly supported at C and D, where 0.5AC = m and 2DB = m, as shown in the figure above. A boy of mass 40 kg stands on the plank at the point M, where M is the midpoint of CD. The plank with the boy standing on the plank, remains in equilibrium with AB horizontal. The plank is modelled as a non-uniform rod and the boy as a particle.
    a) Calculate the magnitude of each of the reaction forces acting on the rod at C and D.
    The boy next moves and stands at the point E on the plank, so that the plank is at the point of tilting about D.
    b) Determine the distance DE.""",

    """Question 12 (***)
    The figure above shows a uniform rod AB of length 4 m and mass 100 kg. The rod rests in equilibrium in a horizontal position, on two supports at C and D, where 0.5AC = m and DB x= m.
    a) Given that the reaction force at the support at D is three times as large as the reaction force at the support at C, determine the value of x.
    The support at D is next moved to a new position E, where 0.75EB = m and an additional mass of m kg is placed at B. The rod remains in equilibrium in a horizontal position and the reaction force at the support at E is now twice as large as the reaction force at the support at C.
    b) Calculate the value of m.""",

    """Question 13 (***+)
    A non-uniform rod AB has length 7 m and weight 300 N. The centre of mass of the rod is x m from A. The rod is placed on two smooth supports at C and D, where 2.5 mAC = and 2 mDB =. The supports at C and D are at the same horizontal level, as shown in the figure above. When a particle of weight W N is placed on the rod at A the reaction force on the rod at C is 200 N. The rod and the particle rest in equilibrium, with AB in a horizontal position.
    a) Show clearly that 200 60x W= −.
    The particle is then removed from A and placed on the rod at B. The rod and the particle remain in equilibrium, with AB in a horizontal position and the reaction force on the rod at C is now 80 N.
    b) Calculate the value of W and the value of x.""",

    """Question 14 (***+)
    A uniform rod AB has length 5 m and weight 300 N. The rod rests in a horizontal position on two smooth supports at C and D, where 1 mAC = and 2 mDB =, as shown in the figure above. A particle of weight W N is placed on the rod at the point E, where AE x= m. The magnitude of the reaction on the rod at C is twice the magnitude of the reaction on the rod at D.
    a) Show clearly that 750/5 3 W x −.
    b) Determine the range of possible values of x.""",

    """Question 15 (****)
    A thin rigid non-uniform beam AB of length 6 m and weight 800 N has its centre of mass at G, where 4AG = m. An additional weight of 100 N is fixed at A. The beam lies in a horizontal position supported by a rough peg at C, where 1AC = m, and a light inextensible wire attached at B. When the wire is inclined at an angle θ to the horizontal, where sin 0.8θ =, the beam remains horizontal, in limiting equilibrium. Calculate the tension in the wire and the value of the coefficient of friction between th...
    """,

    """Question 16 (****)
    A rod AB has mass m kg and length 4 m. The rod is hanging in equilibrium in a horizontal position by two vertical strings attached to the rod. The rod is uniform and the strings are light and inextensible. One string is attached to A and the other string is attached to the point C on the rod as shown in the figure above. The tension in the string attached at C is twice as large as the tension in the string attached at A. Then a particle of mass mλ kg is attached to the rod at B. The rod remains in...
    """,

    """Question 17 (****)
    The standard unit vectors i and j are oriented in the positive x direction and positive y direction, respectively. Three forces 1 4 b= +F i j, 2 3 2a b= +F i j and 3 10 3b= +F i j, where a and b are scalar constants, are acting at the points ( )1,2A, ( )4, 2B − and ( )3, 5C − −, respectively.
    a) Given that the resultant of the three forces is zero, determine the magnitude and direction of the total moment of these three forces about O.
    b) Find, by direct calculation, the magnitude and direction of the total moment of these three forces about C.""",

    """Question 18 (****)
    A non-uniform plank AB has length 12 m and mass M kg. A smooth support is placed under the plank at the point C, where 3AC = m. When a child of mass 30 kg stands at A, the plank rests horizontally in equilibrium. The smooth support is next placed under the plank at the point D, where 5BD = m. When the same child stands at B, the plank again rests horizontally in equilibrium. The plank is modelled as a non-uniform rod whose centre of mass is at the point G, and the child is modelled as a particle.
    a) Determine the value of M.
    b) Calculate the distance AG.
    Two smooth supports are next placed under the plank at the points C and D, and when the same child stands at E, the plank rests horizontally in equilibrium with the reactions at the two supports being equal.
    c) Find the distance AE.""",

    """Question 19 (*****)
    The figure above shows a light rigid framework ABCD, where 90BDA = ° and 90BCD = °.
    It is also given that 0.37AB = m, 0.28BC = m, 0.21CD = m and 0.12AD = m. Forces of magnitude 185 N, 84 N, 63 N and 60 N are acting along AB, BC, CD and AD, in the directions indicated by the arrows in the figure. The 4 forces reduce to a single force acting at some point P on AD, and at right angles to AD. Determine the distance AP.""",

    """Question 20 (*****)
    A uniform rod AB has length 5 m and weight 100 N. The rod rests in a horizontal position on two smooth supports at P and Q, where 1 mAP =, as shown in the figure above. The magnitude of the reaction force on the rod at P is 40 N.
    a) Determine the magnitude of the reaction force on the rod at Q.
    b) Calculate the distance AQ."""
]

# Create slides and add questions to each slide
for question in questions:
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Use a blank slide layout
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5.5))
    text_frame = textbox.text_frame
    p = text_frame.add_paragraph()
    p.text = question

# Save the presentation
output_path = r"C:\Users\Lreps\Desktop\Questions_Presentation.pptx"
prs.save(output_path)
output_path
